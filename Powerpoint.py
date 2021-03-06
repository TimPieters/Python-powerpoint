from pptx import Presentation
import os
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
import pandas as pd
import numpy as np
from PIL import Image

df = pd.read_excel("D:/Upthrust/Frank/PoC Cookbook Automation/Excel/Volkswagen Cookbook Excel.xlsx")

df['Started'] = df['Started'].dt.strftime('%m/%d/%Y')

prs = Presentation("D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Ok.pptx")
#16:9 powerpoint met Upthrust Template

title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "VW Cookbook"
subtitle.text = "Growthmarketing"


for index, row in df.iterrows():

    header = row["Campaigns"]
    started = str(row["Started"])
    status = row["Status"]
    channels = row["Channels"]
    total_reach = str(row["Total Reach"])
    total_clicks = str(row["Total Clicks"])
    media_spend = "€" + str(row["Media Spend"])
    total_leads = str(row["Total Leads"])
    cost_per_lead = "€"+ str(row["Cost Per Lead"])
    ads_set_up = str(row["Ads Set Up"])
    picture_1 = str(row["Picture1"])
    picture_2 = str(row["Picture2"])
    picture_3 = str(row["Picture3"])
    picture_4 = str(row["Picture4"])

    traffic_generation = row["Traffic Generation"]
    landing_page = str(row["Landing Page"])

    gray_border = "D:\\Upthrust\\Frank\\PoC Cookbook Automation\\Excel\\Pictures\\gray_border.png"
    arrow = "D:\\Upthrust\\Frank\\PoC Cookbook Automation\\Excel\\Pictures\\arrow.png"

    def main(i):

        normal_slide = prs.slide_layouts[0]
        slide_2 = prs.slides.add_slide(normal_slide)
        title = slide_2.shapes.title

        img_path = gray_border

        top = Inches(6.5)
        left = Inches(0)
        height = Inches(1.1)
        pic = slide_2.shapes.add_picture(img_path, left, top, height=height)
        #een grijze border, puur estetisch.

        img_path = arrow

        top = Inches(3)
        left = Inches(6.5)
        height = Inches(2)
        pic = slide_2.shapes.add_picture(img_path, left, top, height=height)
        #een rode pijl, puur estetisch.

        if picture_1 != "nan":
            im = Image.open(picture_1)
            width, height = im.size
            image_ratio = width / height
            img_path = picture_1

            if image_ratio < 1.2:
                top = Inches(3)
                left = Inches(0.5)
                height = Inches(2)
            else:
                top = Inches(2)
                left = Inches(0.5)
                height = Inches(2)

            pic = slide_2.shapes.add_picture(img_path, left, top, height=height)

        if picture_2 != "nan":
            im = Image.open(picture_2)
            width, height = im.size
            image_ratio = width / height
            img_path = picture_2

            if image_ratio < 1.2:
                top = Inches(3)
                left = Inches(3)
                height = Inches(2)
            else:
                top = Inches(4)
                left = Inches(0.5)
                height = Inches(2)

            pic = slide_2.shapes.add_picture(img_path, left, top, height=height)

        if picture_3 != "nan":
            im = Image.open(picture_3)
            width, height = im.size
            image_ratio = width / height
            img_path = picture_3

            if image_ratio < 1.2:
                top = Inches(2)
                left = Inches(7.5)
                height = Inches(4)
            else:
                top = Inches(2)
                left = Inches(7.5)
                height = Inches(4)

            pic = slide_2.shapes.add_picture(img_path, left, top, height=height)

        if picture_4 != "nan":
            im = Image.open(picture_4)
            width, height = im.size
            image_ratio = width / height
            img_path = picture_4

            if image_ratio < 1.2:
                top = Inches(3.5)
                left = Inches(10)
                height = Inches(1.2)
            else:
                top = Inches(3.5)
                left = Inches(10)
                height = Inches(1.2)

            pic = slide_2.shapes.add_picture(img_path, left, top, height=height)

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame       # do things with the text frame

        text_frame.clear()
        p = title.text_frame.paragraphs[0]
        run = p.add_run()

        font = run.font
        font.name = 'Century Goth'
        font.size = Pt(25.3)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        run.text = header

        def bold_text():
            font = run.font
            font.name = 'Century Goth'
            font.size = Pt(14)
            font.bold = True
            font.italic = None  # cause value to be inherited from theme

        def normal_text():
            font = run.font
            font.name = 'Century Goth'
            font.size = Pt(14)
            font.bold = False
            font.italic = None  # cause value to be inherited from theme

        def small_text():
            font = run.font
            font.name = 'Century Goth'
            font.size = Pt(12)
            font.bold = False
            font.italic = None  # cause value to be inherited from theme

        txBox = slide_2.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(2), Inches(1))
        tf = txBox.text_frame

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue                    # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        bold_text()
        run.text = 'Started: '

        run = p.add_run()
        normal_text()
        run.text = started

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Status: '

        run = p.add_run()
        normal_text()
        run.text = status

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Channel(s): '

        run = p.add_run()
        normal_text()
        run.text = channels

        txBox = slide_2.shapes.add_textbox(Inches(4), Inches(6.5), Inches(2), Inches(1))
        tf = txBox.text_frame

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue                # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        bold_text()
        run.text = 'Current Results: '

        txBox = slide_2.shapes.add_textbox(Inches(7), Inches(6.5), Inches(2), Inches(1))
        tf = txBox.text_frame

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue                # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        bold_text()
        run.text = 'Total Reach: '

        run = p.add_run()
        normal_text()
        run.text = total_reach

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Total Clicks: '

        run = p.add_run()
        normal_text()
        run.text = total_clicks

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Media Spend: '

        run = p.add_run()
        normal_text()
        run.text = media_spend

        txBox = slide_2.shapes.add_textbox(Inches(10), Inches(6.5), Inches(2), Inches(1))
        tf = txBox.text_frame


        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue                # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        bold_text()
        run.text = 'Total Leads: '

        run = p.add_run()
        normal_text()
        run.text = total_leads

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Cost Per Lead: : '

        run = p.add_run()
        normal_text()
        run.text = cost_per_lead

        p = tf.add_paragraph()

        run = p.add_run()
        bold_text()
        run.text = 'Ads Set Up: : '

        run = p.add_run()
        normal_text()
        run.text = ads_set_up

        txBox = slide_2.shapes.add_textbox(Inches(1.6), Inches(1.8), Inches(2), Inches(0.3))
        tf = txBox.text_frame

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue  # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        small_text()
        run.text = "Traffic generation: "

        if traffic_generation != "nan":

            run = p.add_run()
            small_text()
            run.text = traffic_generation

        txBox = slide_2.shapes.add_textbox(Inches(8.5), Inches(1.7), Inches(4), Inches(0.3))
        tf = txBox.text_frame

        for shape in slide_2.shapes:
            if not shape.has_text_frame:
                continue  # do things with the text frame

        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]


        run = p.add_run()
        small_text()
        run.text = "Landing page: "

        if landing_page != "nan":

            run = p.add_run()
            small_text()
            run.text = landing_page
            run.hyperlink.address = landing_page
        prs.save('D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx')

    i = 0
    if i <= index:
        main(i)
        i += 1


os.startfile("D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx")









